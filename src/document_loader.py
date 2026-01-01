"""
Módulo para cargar diferentes tipos de documentos
"""
from pathlib import Path
from typing import List, Dict
import json
import pypdf
import openpyxl
import xlrd
import markdown
import ebooklib
from ebooklib import epub  
from docx import Document
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from rich.console import Console
from config.settings import SUPPORTED_EXTENSIONS

console = Console()


class DocumentLoader:
    """Carga documentos de diferentes formatos"""
    
    def __init__(self):
        self.supported_extensions = SUPPORTED_EXTENSIONS
    
    def load_document(self, file_path: str) -> Dict[str, any]:
        """
        Carga un documento y retorna su contenido
        
        Args:
            file_path: Ruta al archivo
            
        Returns:
            Dict con 'content', 'metadata', 'file_type'
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"Archivo no encontrado: {file_path}")
        
        extension = path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Formato no soportado: {extension}")
        
        console.print(f"[blue]Cargando:[/blue] {path.name}")
        
        file_type = self.supported_extensions[extension]
        
        # Mapeo de tipos a métodos
        loaders = {
            'text': self._load_txt,
            'markdown': self._load_markdown,
            'pdf': self._load_pdf,
            'docx': self._load_docx,
            'doc': self._load_docx,
            'excel': self._load_excel,
            'excel_old': self._load_excel_old,
            'powerpoint': self._load_pptx,
            'powerpoint_old': self._load_ppt_placeholder,
            'csv': self._load_csv,
            'json': self._load_json,
            'xml': self._load_xml,
            'html': self._load_html,
            'rtf': self._load_rtf,
            'epub': self._load_epub,
        }
        
        loader = loaders.get(file_type)
        if loader:
            return loader(path)
        else:
            raise ValueError(f"No hay loader para el tipo: {file_type}")
    
    def _load_txt(self, path: Path) -> Dict:
        """Carga archivo de texto plano"""
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            'content': content,
            'metadata': {
                'source': str(path),
                'file_name': path.name,
                'file_type': 'text'
            },
            'file_type': 'text'
        }
    
    def _load_markdown(self, path: Path) -> Dict:
        """Carga archivo Markdown y lo convierte a texto"""
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            md_content = f.read()
        
        # Convertir Markdown a HTML y luego extraer texto
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        text_content = soup.get_text(separator='\n\n')
        
        return {
            'content': text_content,
            'metadata': {
                'source': str(path),
                'file_name': path.name,
                'file_type': 'markdown'
            },
            'file_type': 'markdown'
        }
    
    def _load_pdf(self, path: Path) -> Dict:
        """Carga archivo PDF y detecta si necesita OCR"""
        content = []
        
        try:
            with open(path, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        content.append(text)
            
            full_content = '\n\n'.join(content)
            
            # Verificar si el PDF tiene contenido útil
            if not full_content.strip() or len(full_content.strip()) < 50:
                console.print(f"[yellow]⚠ ADVERTENCIA:[/yellow] {path.name}")
                console.print("[yellow]  → Este PDF parece estar escaneado o sin texto extraíble[/yellow]")
                console.print("[yellow]  → Se necesita OCR para procesarlo[/yellow]")
                console.print("[yellow]  → Por ahora será ignorado[/yellow]\n")
                
                return {
                    'content': '',
                    'metadata': {
                        'source': str(path),
                        'file_name': path.name,
                        'file_type': 'pdf',
                        'num_pages': num_pages,
                        'needs_ocr': True,
                        'warning': 'PDF sin texto extraíble'
                    },
                    'file_type': 'pdf'
                }
            
            return {
                'content': full_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'pdf',
                    'num_pages': num_pages,
                    'needs_ocr': False
                },
                'file_type': 'pdf'
            }
        except Exception as e:
            console.print(f"[red]Error al leer PDF:[/red] {e}")
            raise
    
    def _load_docx(self, path: Path) -> Dict:
        """Carga archivo Word"""
        try:
            doc = Document(path)
            content = []
            
            # Extraer párrafos
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            # Extraer texto de tablas
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        content.append(row_text)
            
            full_content = '\n\n'.join(content)
            
            return {
                'content': full_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'docx',
                    'num_paragraphs': len(doc.paragraphs),
                    'num_tables': len(doc.tables)
                },
                'file_type': 'docx'
            }
        except Exception as e:
            console.print(f"[red]Error al leer Word:[/red] {e}")
            raise
    
    def _load_excel(self, path: Path) -> Dict:
        """Carga archivo Excel moderno (.xlsx)"""
        try:
            # Leer todas las hojas
            excel_file = pd.ExcelFile(path)
            content = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                
                # Convertir DataFrame a texto estructurado
                sheet_text = f"--- Hoja: {sheet_name} ---\n"
                sheet_text += df.to_string(index=False)
                content.append(sheet_text)
            
            full_content = '\n\n'.join(content)
            
            return {
                'content': full_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'excel',
                    'num_sheets': len(excel_file.sheet_names),
                    'sheet_names': excel_file.sheet_names
                },
                'file_type': 'excel'
            }
        except Exception as e:
            console.print(f"[red]Error al leer Excel:[/red] {e}")
            raise
    
    def _load_excel_old(self, path: Path) -> Dict:
        """Carga archivo Excel antiguo (.xls)"""
        try:
            workbook = xlrd.open_workbook(path)
            content = []
            
            for sheet in workbook.sheets():
                sheet_text = f"--- Hoja: {sheet.name} ---\n"
                rows = []
                for row_idx in range(sheet.nrows):
                    row = [str(cell.value) for cell in sheet.row(row_idx)]
                    rows.append(' | '.join(row))
                sheet_text += '\n'.join(rows)
                content.append(sheet_text)
            
            full_content = '\n\n'.join(content)
            
            return {
                'content': full_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'excel_old',
                    'num_sheets': workbook.nsheets
                },
                'file_type': 'excel'
            }
        except Exception as e:
            console.print(f"[red]Error al leer Excel antiguo:[/red] {e}")
            raise
    
    def _load_pptx(self, path: Path) -> Dict:
        """Carga archivo PowerPoint (.pptx)"""
        try:
            prs = Presentation(path)
            content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Diapositiva {i} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + '\n'
                
                if slide_text.strip() != f"--- Diapositiva {i} ---":
                    content.append(slide_text)
            
            full_content = '\n\n'.join(content)
            
            return {
                'content': full_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'powerpoint',
                    'num_slides': len(prs.slides)
                },
                'file_type': 'powerpoint'
            }
        except Exception as e:
            console.print(f"[red]Error al leer PowerPoint:[/red] {e}")
            raise
    
    def _load_ppt_placeholder(self, path: Path) -> Dict:
        """Placeholder para archivos .ppt antiguos"""
        console.print(f"[yellow]⚠ Formato .ppt no soportado directamente[/yellow]")
        console.print(f"[yellow]  → Convierte a .pptx para mejor compatibilidad[/yellow]")
        return {
            'content': '',
            'metadata': {
                'source': str(path),
                'file_name': path.name,
                'file_type': 'powerpoint_old',
                'warning': 'Formato .ppt antiguo no soportado'
            },
            'file_type': 'powerpoint'
        }
    
    def _load_csv(self, path: Path) -> Dict:
        """Carga archivo CSV"""
        try:
            df = pd.read_csv(path, encoding='utf-8', on_bad_lines='skip')
            
            # Convertir a texto estructurado
            content = df.to_string(index=False)
            
            return {
                'content': content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'csv',
                    'num_rows': len(df),
                    'num_columns': len(df.columns),
                    'columns': list(df.columns)
                },
                'file_type': 'csv'
            }
        except Exception as e:
            console.print(f"[red]Error al leer CSV:[/red] {e}")
            raise
    
    def _load_json(self, path: Path) -> Dict:
        """Carga archivo JSON"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convertir JSON a texto legible
            content = json.dumps(data, indent=2, ensure_ascii=False)
            
            return {
                'content': content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'json'
                },
                'file_type': 'json'
            }
        except Exception as e:
            console.print(f"[red]Error al leer JSON:[/red] {e}")
            raise
    
    def _load_xml(self, path: Path) -> Dict:
        """Carga archivo XML"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                xml_content = f.read()
            
            soup = BeautifulSoup(xml_content, 'lxml-xml')
            text_content = soup.get_text(separator='\n', strip=True)
            
            return {
                'content': text_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'xml'
                },
                'file_type': 'xml'
            }
        except Exception as e:
            console.print(f"[red]Error al leer XML:[/red] {e}")
            raise
    
    def _load_html(self, path: Path) -> Dict:
        """Carga archivo HTML"""
        try:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Eliminar scripts y estilos
            for script in soup(["script", "style"]):
                script.decompose()
            
            text_content = soup.get_text(separator='\n\n', strip=True)
            
            return {
                'content': text_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'html'
                },
                'file_type': 'html'
            }
        except Exception as e:
            console.print(f"[red]Error al leer HTML:[/red] {e}")
            raise
    
    def _load_rtf(self, path: Path) -> Dict:
        """Carga archivo RTF"""
        try:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text_content = rtf_to_text(rtf_content)
            
            return {
                'content': text_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'rtf'
                },
                'file_type': 'rtf'
            }
        except Exception as e:
            console.print(f"[red]Error al leer RTF:[/red] {e}")
            raise
    def _load_epub(self, path: Path) -> Dict:
        """Carga archivo EPUB (libro electrónico)"""
        try:
            book = epub.read_epub(str(path))
            content = []
            
            # Extraer metadata del libro
            title = book.get_metadata('DC', 'title')
            title = title[0][0] if title else 'Sin título'
            
            author = book.get_metadata('DC', 'creator')
            author = author[0][0] if author else 'Desconocido'
            
            # Extraer texto de todos los capítulos
            for item in book.get_items():
                # Solo procesar documentos HTML (capítulos)
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Parsear HTML del capítulo
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    
                    # Eliminar scripts y estilos
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    # Extraer texto limpio
                    chapter_text = soup.get_text(separator='\n', strip=True)
                    
                    if chapter_text.strip():
                        content.append(chapter_text)
            
            full_content = '\n\n'.join(content)
            
            # Verificar que haya contenido
            if not full_content.strip():
                console.print(f"[yellow]⚠ ADVERTENCIA:[/yellow] {path.name}")
                console.print("[yellow]  → Este EPUB no tiene texto extraíble[/yellow]\n")
                
                return {
                    'content': '',
                    'metadata': {
                        'source': str(path),
                        'file_name': path.name,
                        'file_type': 'epub',
                        'title': title,
                        'author': author,
                        'warning': 'EPUB sin contenido'
                    },
                    'file_type': 'epub'
                }
            
            return {
                'content': full_content,
                'metadata': {
                    'source': str(path),
                    'file_name': path.name,
                    'file_type': 'epub',
                    'title': title,
                    'author': author,
                    'num_chapters': len(content)
                },
                'file_type': 'epub'
            }
        except Exception as e:
            console.print(f"[red]Error al leer EPUB:[/red] {e}")
            raise
    
    def load_directory(self, directory_path: str) -> List[Dict]:
        """
        Carga todos los documentos soportados de un directorio
        
        Args:
            directory_path: Ruta al directorio
            
        Returns:
            Lista de documentos cargados
        """
        path = Path(directory_path)
        documents = []
        skipped = 0
        
        if not path.is_dir():
            raise ValueError(f"No es un directorio válido: {directory_path}")
        
        for file_path in path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                try:
                    doc = self.load_document(str(file_path))
                    
                    # Solo agregar si tiene contenido
                    if doc['content'].strip():
                        documents.append(doc)
                    else:
                        skipped += 1
                        
                except Exception as e:
                    console.print(f"[yellow]Advertencia:[/yellow] No se pudo cargar {file_path.name}: {e}")
                    skipped += 1
        
        if skipped > 0:
            console.print(f"[yellow]⚠ {skipped} archivo(s) omitido(s) (sin contenido o con errores)[/yellow]")
        
        console.print(f"[green]✓[/green] Cargados {len(documents)} documentos")
        return documents